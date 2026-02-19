"""Utilities for processing different document types prior to embedding.

The module exposes a simple factory that maps MIME content types (or file
extensions) to dedicated processors. Each processor is responsible for
transforming the raw file bytes into semantic chunks that downstream
components (e.g., embeddings, vector stores) can consume.
"""

from __future__ import annotations

import base64
import gc
import io
import logging
import os
import time
import xml.etree.ElementTree as ET
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple, Type

import boto3
import pdfplumber
import psutil
from botocore.config import Config
from botocore.exceptions import ClientError
from docx import Document as DocxDocument
from docx.oxml.ns import qn
from docx.table import Table
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.config.document_processing_config import (
    IMAGE_ANALYSIS_CONFIG,
    MAX_IMAGE_SIZE_BYTES,
    SUPPORTED_IMAGE_FORMATS,
)
from backend.config.prompts import IMAGE_ANALYSIS_PROMPT

logger = logging.getLogger(__name__)

_process = psutil.Process(os.getpid())


def _log_memory(tag: str) -> None:
    """Log current process RSS memory at INFO level."""
    try:
        rss_mb = _process.memory_info().rss / (1024 * 1024)
        logger.info("[MEM] %s: RSS=%.1f MB", tag, rss_mb)
    except Exception:
        pass


Chunk = Dict[str, object]
Chunker = Callable[[Dict[str, object], str], list[Chunk]]


def has_cid_corruption(text: str) -> bool:
    """
    Check if extracted text contains CID encoding corruption.
    
    CID (Character Identifier) corruption appears as patterns like:
    (cid:6)(cid:7)(cid:1) etc., indicating symbolic fonts that can't be decoded.
    
    Args:
        text: Extracted text to check
    
    Returns:
        True if CID corruption is detected
    """
    return "(cid:" in text.lower()


class DocumentProcessor(ABC):
    """Base interface for document processors.
    
    Provides shared functionality for image analysis via AWS Bedrock,
    which can be used by all subclasses when processing documents with images.
    """

    _bedrock_client = None

    @classmethod
    def _get_bedrock_client(cls):
        """Get or create a shared Bedrock Runtime client."""
        if cls._bedrock_client is None:
            config = Config(
                region_name=IMAGE_ANALYSIS_CONFIG.get("region", "eu-central-1"),
                retries={
                    "max_attempts": IMAGE_ANALYSIS_CONFIG.get("max_retries", 3),
                    "mode": "adaptive",
                },
            )
            cls._bedrock_client = boto3.client("bedrock-runtime", config=config)
        return cls._bedrock_client

    def analyze_image(
        self,
        image_bytes: bytes,
        image_format: str,
        context: str = "",
    ) -> Optional[str]:
        """Analyze an image using AWS Bedrock Nova Lite vision model.
        
        Args:
            image_bytes: Raw bytes of the image
            image_format: Format of the image (png, jpeg, gif, webp)
            context: Optional context about where the image appears in the document
            
        Returns:
            Description of the image including any text found, or None if analysis fails
        """
        # Validate image size
        if len(image_bytes) > MAX_IMAGE_SIZE_BYTES:
            logger.warning(
                "Image exceeds maximum size (%d bytes > %d bytes), skipping analysis",
                len(image_bytes),
                MAX_IMAGE_SIZE_BYTES,
            )
            return None

        # Validate format
        if image_format.lower() not in ["png", "jpeg", "gif", "webp"]:
            logger.warning("Unsupported image format: %s", image_format)
            return None

        # Build the prompt with optional context
        prompt = IMAGE_ANALYSIS_PROMPT
        if context:
            prompt = f"{context}\n\n{prompt}"

        # Build the message with image content
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "image": {
                            "format": image_format.lower(),
                            "source": {
                                "bytes": image_bytes,
                            },
                        },
                    },
                    {
                        "text": prompt,
                    },
                ],
            }
        ]

        try:
            return self._call_bedrock_with_retry(messages)
        except Exception as e:
            logger.error("Image analysis failed: %s", str(e))
            return None

    def analyze_multiple_images(
        self,
        images: List[Tuple[bytes, str]],
        context: str = "",
    ) -> Optional[str]:
        """Analyze multiple images in a single Bedrock call.
        
        Args:
            images: List of tuples (image_bytes, image_format)
            context: Optional context about where the images appear
            
        Returns:
            Combined description of all images, or None if analysis fails
        """
        if not images:
            return None

        # Filter valid images
        valid_images = []
        for image_bytes, image_format in images:
            if len(image_bytes) > MAX_IMAGE_SIZE_BYTES:
                logger.warning("Skipping oversized image in batch")
                continue
            if image_format.lower() not in ["png", "jpeg", "gif", "webp"]:
                logger.warning("Skipping unsupported format in batch: %s", image_format)
                continue
            valid_images.append((image_bytes, image_format))

        if not valid_images:
            return None

        # Build content blocks with all images
        content_blocks: List[dict] = []
        for image_bytes, image_format in valid_images:
            content_blocks.append({
                "image": {
                    "format": image_format.lower(),
                    "source": {
                        "bytes": image_bytes,
                    },
                },
            })

        # Add the prompt
        prompt = IMAGE_ANALYSIS_PROMPT
        if context:
            prompt = f"{context}\n\n{prompt}"
        content_blocks.append({"text": prompt})

        messages = [{"role": "user", "content": content_blocks}]

        try:
            return self._call_bedrock_with_retry(messages)
        except Exception as e:
            logger.error("Multi-image analysis failed: %s", str(e))
            return None

    def _call_bedrock_with_retry(self, messages: List[dict]) -> str:
        """Call Bedrock API with retry logic for throttling.
        
        Args:
            messages: The messages to send to the model
            
        Returns:
            The model's response text
            
        Raises:
            Exception: If the API call fails after retries
        """
        client = self._get_bedrock_client()
        model_id = IMAGE_ANALYSIS_CONFIG["model_id"]
        max_retries = IMAGE_ANALYSIS_CONFIG.get("max_retries", 3)
        retry_delay = IMAGE_ANALYSIS_CONFIG.get("retry_delay_seconds", 5)

        for attempt in range(max_retries):
            try:
                response = client.converse(
                    modelId=model_id,
                    messages=messages,
                    inferenceConfig={
                        "temperature": IMAGE_ANALYSIS_CONFIG.get("temperature", 0.3),
                        "maxTokens": IMAGE_ANALYSIS_CONFIG.get("max_tokens", 2000),
                        "topP": IMAGE_ANALYSIS_CONFIG.get("top_p", 0.9),
                    },
                )

                # Extract the response text
                output_message = response.get("output", {}).get("message", {})
                content_blocks = output_message.get("content", [])

                response_text = ""
                for block in content_blocks:
                    if "text" in block:
                        response_text += block["text"]

                # Log the full response text
                # logger.info("Image analysis response text: %s", response_text)

                # Log token usage
                usage = response.get("usage", {})
                logger.debug(
                    "Image analysis tokens - Input: %d, Output: %d",
                    usage.get("inputTokens", 0),
                    usage.get("outputTokens", 0),
                )

                return response_text

            except ClientError as exc:
                error_code = (exc.response.get("Error") or {}).get("Code")
                if error_code == "ThrottlingException" and attempt < max_retries - 1:
                    logger.warning(
                        "Bedrock throttled, sleeping %d seconds (attempt %d/%d)",
                        retry_delay,
                        attempt + 1,
                        max_retries,
                    )
                    time.sleep(retry_delay)
                    retry_delay *= 2  # Exponential backoff
                    continue
                raise

        raise RuntimeError("Bedrock API call failed after all retries")

    @abstractmethod
    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        """Return a list of chunks extracted from the document."""


class PDFDocumentProcessor(DocumentProcessor):
    """Process PDF documents into text chunks.
    
    Skips pages with CID encoding corruption (symbolic fonts).
    Extracts and analyzes images using Bedrock vision model.
    Falls back to page-as-image OCR when no text is extractable.
    """

    def _render_page_as_image(self, page) -> Optional[Tuple[bytes, str]]:
        """Render an entire PDF page as a PNG image for OCR/analysis.
        
        Used as fallback when text extraction fails (e.g., scanned PDFs).
        Optimized for Cloud Run with memory management and size limits.
        
        Args:
            page: A pdfplumber page object
            
        Returns:
            Tuple of (image_bytes, format) or None if rendering fails
        """
        try:
            # Render the page at 150 DPI (good balance of quality and size for Cloud Run)
            pil_image = page.to_image(resolution=150).original
            
            # Convert to PNG bytes
            img_buffer = io.BytesIO()
            pil_image.save(img_buffer, format="PNG")
            image_bytes = img_buffer.getvalue()
            
            # Clean up immediately to free memory
            del pil_image
            del img_buffer
            
            # Check size and retry at lower resolution if needed
            if len(image_bytes) > MAX_IMAGE_SIZE_BYTES:
                logger.warning(
                    "Rendered page exceeds max size (%d > %d), trying lower resolution",
                    len(image_bytes),
                    MAX_IMAGE_SIZE_BYTES
                )
                # Try again at 100 DPI
                pil_image = page.to_image(resolution=100).original
                img_buffer = io.BytesIO()
                pil_image.save(img_buffer, format="PNG")
                image_bytes = img_buffer.getvalue()
                del pil_image
                del img_buffer
                
                if len(image_bytes) > MAX_IMAGE_SIZE_BYTES:
                    logger.warning("Page still too large at 100 DPI, skipping")
                    return None
            
            return (image_bytes, "png")
        except Exception as e:
            logger.debug("Failed to render PDF page as image: %s", e)
            return None

    def _extract_page_images(self, page) -> List[Tuple[bytes, str]]:
        """Extract all images from a PDF page.
        
        Uses pdfplumber to crop and render each image region as PNG,
        ensuring the output is a valid image format for Bedrock.
        
        Args:
            page: A pdfplumber page object
            
        Returns:
            List of tuples (image_bytes, format)
        """
        images = []
        try:
            # pdfplumber provides image metadata including bounding boxes
            for img in page.images:
                try:
                    # Get image bounding box
                    x0 = img.get("x0", 0)
                    top = img.get("top", 0)
                    x1 = img.get("x1", 0)
                    bottom = img.get("bottom", 0)
                    
                    # Skip if bounding box is invalid or too small
                    if x1 <= x0 or bottom <= top:
                        continue
                    if (x1 - x0) < 10 or (bottom - top) < 10:
                        continue
                    
                    # Crop the page to the image region and convert to PIL Image
                    cropped = page.crop((x0, top, x1, bottom))
                    pil_image = cropped.to_image(resolution=150).original
                    # Free the cropped page object immediately
                    cropped.flush_cache()
                    del cropped
                    
                    # Convert to PNG bytes
                    img_buffer = io.BytesIO()
                    pil_image.save(img_buffer, format="PNG")
                    del pil_image
                    image_bytes = img_buffer.getvalue()
                    del img_buffer
                    
                    if image_bytes and len(image_bytes) > 100:  # Skip tiny/empty images
                        images.append((image_bytes, "png"))
                    del image_bytes
                        
                except Exception as e:
                    logger.debug("Failed to extract image from PDF page: %s", e)
                    continue
        except Exception as e:
            logger.debug("Failed to iterate images on PDF page: %s", e)
        
        return images

    # Number of pages processed per pdfplumber context.
    # After each batch the PDF object is closed, releasing all pdfminer internal
    # caches (XRef tables, object streams, decoded content) that accumulate ~2-3 MB/page
    # and cannot be freed any other way while the PDF context remains open.
    _PAGE_BATCH_SIZE = 20

    def _process_page(
        self,
        page,
        page_number: int,
        total_pages: int,
        doc_name: str,
        storage_path: str,
        chunk_document: Chunker,
    ) -> list[Chunk]:
        """Process a single PDF page and return its chunks. Frees all per-page objects."""
        page_chunks: list[Chunk] = []

        page_text = page.extract_text() or ""

        # Skip pages with CID corruption (symbolic fonts that can't be decoded)
        if page_text.strip() and has_cid_corruption(page_text):
            logger.warning(
                "[PDF] Skipping page %d of '%s': CID encoding corruption detected",
                page_number, doc_name,
            )
            del page_text
            page.flush_cache()
            gc.collect()
            _log_memory(f"[PDF] page {page_number}/{total_pages} (CID skip) '{doc_name}'")
            return page_chunks

        # Extract and analyze images on this page
        page_images = self._extract_page_images(page)
        has_page_images = len(page_images) > 0
        image_descriptions: List[str] = []

        if page_images:
            logger.info(
                "[PDF] Found %d image(s) on page %d of '%s'",
                len(page_images), page_number, doc_name,
            )
            if len(page_images) == 1:
                description = self.analyze_image(
                    page_images[0][0],
                    page_images[0][1],
                    context=f"Image from page {page_number} of PDF document '{doc_name}'",
                )
                if description:
                    image_descriptions.append(description)
            else:
                description = self.analyze_multiple_images(
                    page_images,
                    context=f"Images from page {page_number} of PDF document '{doc_name}'",
                )
                if description:
                    image_descriptions.append(description)

        # Free image bytes immediately after analysis
        del page_images

        # Release pdfplumber's page-level cache
        page.flush_cache()

        # Combine text and image descriptions
        combined_content_parts: List[str] = []

        if page_text.strip():
            combined_content_parts.append(page_text.strip())
        del page_text

        if image_descriptions:
            combined_content_parts.append("\n[Image Content]\n" + "\n".join(image_descriptions))
        del image_descriptions

        if not combined_content_parts:
            gc.collect()
            _log_memory(f"[PDF] page {page_number}/{total_pages} (empty) '{doc_name}'")
            return page_chunks

        combined_text = "\n\n".join(combined_content_parts)
        del combined_content_parts

        doc_metadata = {
            "name": doc_name,
            "page": page_number,
            "storage_path": storage_path,
            "has_images": has_page_images,
        }
        page_chunks.extend(chunk_document(doc_metadata, combined_text))
        del combined_text

        gc.collect()
        _log_memory(f"[PDF] page {page_number}/{total_pages} '{doc_name}'")
        return page_chunks

    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        chunks: list[Chunk] = []

        # Quick first open just to read the page count, then close immediately.
        with pdfplumber.open(io.BytesIO(file_bytes)) as _pdf:
            total_pages = len(_pdf.pages)

        # Skip very large PDFs to avoid OOM
        if total_pages > 400:
            logger.warning(
                "[PDF] Skipping '%s': %d pages exceeds limit of 400 pages",
                doc_name, total_pages,
            )
            return chunks

        # First pass: text + embedded-image extraction, processed in batches.
        # Each batch opens a fresh pdfplumber (and thus pdfminer) context so that
        # pdfminer's document-level caches (~2-3 MB/page) are fully released between
        # batches rather than accumulating for the entire document.
        for batch_start in range(0, total_pages, self._PAGE_BATCH_SIZE):
            batch_end = min(batch_start + self._PAGE_BATCH_SIZE, total_pages)
            logger.debug(
                "[PDF] Opening batch pages %d-%d/%d of '%s'",
                batch_start + 1, batch_end, total_pages, doc_name,
            )

            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page_number in range(batch_start + 1, batch_end + 1):
                    page = pdf.pages[page_number - 1]
                    page_chunks = self._process_page(
                        page, page_number, total_pages, doc_name, storage_path, chunk_document,
                    )
                    chunks.extend(page_chunks)
            # pdfplumber context is now closed → entire pdfminer PDFDocument object
            # is released; force a GC pass to reclaim it promptly.
            gc.collect()
            _log_memory(
                f"[PDF] after batch {batch_start + 1}-{batch_end}/{total_pages} '{doc_name}'"
            )

        # Fallback: if no chunks generated, treat entire pages as images (e.g., scanned PDFs)
        if not chunks:
            logger.info(
                "[PDF] No text/images extracted from '%s', using page-as-image fallback with Bedrock OCR",
                doc_name,
            )

            max_fallback_pages = min(total_pages, 15)
            if total_pages > max_fallback_pages:
                logger.warning(
                    "[PDF] Limiting fallback processing to first %d of %d pages",
                    max_fallback_pages, total_pages,
                )

            # Open one page at a time so pdfminer state never accumulates
            for page_number in range(1, max_fallback_pages + 1):
                try:
                    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                        page = pdf.pages[page_number - 1]
                        page_image = self._render_page_as_image(page)
                        page.flush_cache()

                    # pdfminer context released here
                    gc.collect()

                    if not page_image:
                        logger.warning("[PDF] Failed to render page %d as image", page_number)
                        continue

                    image_bytes, image_format = page_image
                    description = self.analyze_image(
                        image_bytes,
                        image_format,
                        context=f"Full page {page_number} of scanned PDF document '{doc_name}'",
                    )
                    del image_bytes
                    del page_image

                    if not description or not description.strip():
                        logger.debug("[PDF] No content extracted from page %d image", page_number)
                        continue

                    doc_metadata = {
                        "name": doc_name,
                        "page": page_number,
                        "storage_path": storage_path,
                        "has_images": True,
                        "ocr_fallback": True,
                    }
                    chunks.extend(chunk_document(doc_metadata, description))
                    gc.collect()
                    _log_memory(f"[PDF] fallback page {page_number}/{max_fallback_pages} '{doc_name}'")

                except Exception as e:
                    logger.error("[PDF] Error processing page %d as image: %s", page_number, e)
                    continue

            logger.info(
                "[PDF] Fallback complete: extracted %d chunks from %d pages of '%s'",
                len(chunks), max_fallback_pages, doc_name,
            )

        return chunks


class TextDocumentProcessor(DocumentProcessor):
    """Process plain text documents into chunks."""
    
    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        chunks: list[Chunk] = []
        
        # Decode text file (try UTF-8 first, fallback to latin-1)
        try:
            text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            text = file_bytes.decode('latin-1')
        
        if not text.strip():
            logger.info("[TXT] Skipping empty document '%s'", doc_name)
            return chunks
        
        # Text files don't have pages, so we treat the whole file as page 1
        doc_metadata = {
            "name": doc_name,
            "page": 1,
            "storage_path": storage_path,
        }
        chunks.extend(chunk_document(doc_metadata, text))
        return chunks


class DocxDocumentProcessor(DocumentProcessor):
    """Process DOCX documents into text chunks.
    
    Extracts text from paragraphs, tables, headers, and footers.
    Extracts and analyzes embedded images using Bedrock vision model.
    """

    def _get_images_from_paragraph(self, paragraph, doc) -> List[Tuple[bytes, str]]:
        """Extract images from a paragraph.
        
        Args:
            paragraph: A python-docx paragraph object
            doc: The parent Document object (needed to access relationships)
            
        Returns:
            List of tuples (image_bytes, format)
        """
        images = []
        
        # Find drawing elements (inline images)
        drawing_elements = paragraph._element.findall('.//' + qn('w:drawing'))
        for drawing in drawing_elements:
            # Look for blip elements that contain the image reference
            blip_elements = drawing.findall('.//' + qn('a:blip'))
            for blip in blip_elements:
                embed_id = blip.get(qn('r:embed'))
                if embed_id:
                    try:
                        # Get the image part from relationships
                        image_part = doc.part.related_parts.get(embed_id)
                        if image_part:
                            image_bytes = image_part.blob
                            # Determine format from content type
                            content_type = image_part.content_type
                            img_format = SUPPORTED_IMAGE_FORMATS.get(content_type, "png")
                            images.append((image_bytes, img_format))
                    except Exception as e:
                        logger.debug("Failed to extract image from DOCX paragraph: %s", e)
        
        # Also check for legacy picture elements
        picture_elements = paragraph._element.findall('.//' + qn('w:pict'))
        for pict in picture_elements:
            # Legacy images use v:imagedata
            imagedata_elements = pict.findall('.//' + '{urn:schemas-microsoft-com:vml}imagedata')
            for imagedata in imagedata_elements:
                rel_id = imagedata.get(qn('r:id'))
                if rel_id:
                    try:
                        image_part = doc.part.related_parts.get(rel_id)
                        if image_part:
                            image_bytes = image_part.blob
                            content_type = image_part.content_type
                            img_format = SUPPORTED_IMAGE_FORMATS.get(content_type, "png")
                            images.append((image_bytes, img_format))
                    except Exception as e:
                        logger.debug("Failed to extract legacy image from DOCX: %s", e)
        
        return images

    def _has_images_in_paragraph(self, paragraph) -> bool:
        """Check if a paragraph contains embedded images."""
        drawing_elements = paragraph._element.findall('.//' + qn('w:drawing'))
        picture_elements = paragraph._element.findall('.//' + qn('w:pict'))
        return len(drawing_elements) > 0 or len(picture_elements) > 0

    def _extract_table_text(self, table: Table) -> str:
        """Extract text from a table, preserving structure."""
        table_text_parts: List[str] = []
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            table_text_parts.append(" | ".join(row_cells))
        return "\n".join(table_text_parts)

    def _get_images_from_table(self, table: Table, doc) -> List[Tuple[bytes, str]]:
        """Extract all images from a table."""
        images = []
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    images.extend(self._get_images_from_paragraph(paragraph, doc))
        return images

    def _has_images_in_table(self, table: Table) -> bool:
        """Check if a table contains any images."""
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if self._has_images_in_paragraph(paragraph):
                        return True
        return False

    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        chunks: list[Chunk] = []
        
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
        except Exception as e:
            logger.error("[DOCX] Failed to open document '%s': %s", doc_name, e)
            raise ValueError(f"Failed to open DOCX document '{doc_name}': {e}")

        all_text_parts: List[str] = []
        all_images: List[Tuple[bytes, str]] = []
        
        # Process paragraphs
        for paragraph in doc.paragraphs:
            # Extract text
            text = paragraph.text.strip()
            if text:
                all_text_parts.append(text)
            
            # Extract images
            paragraph_images = self._get_images_from_paragraph(paragraph, doc)
            if paragraph_images:
                all_images.extend(paragraph_images)
        
        # Process tables
        for table in doc.tables:
            # Extract text
            table_text = self._extract_table_text(table)
            if table_text.strip():
                all_text_parts.append(table_text)
            
            # Extract images from table
            table_images = self._get_images_from_table(table, doc)
            if table_images:
                all_images.extend(table_images)
        
        # Analyze images if present
        image_descriptions: List[str] = []
        if all_images:
            logger.info("[DOCX] Found %d image(s) in '%s'", len(all_images), doc_name)
            
            if len(all_images) == 1:
                description = self.analyze_image(
                    all_images[0][0],
                    all_images[0][1],
                    context=f"Image from DOCX document '{doc_name}'"
                )
                if description:
                    image_descriptions.append(description)
            else:
                # Batch analyze images (up to 5 at a time to avoid token limits)
                batch_size = 5
                for i in range(0, len(all_images), batch_size):
                    batch = all_images[i:i + batch_size]
                    description = self.analyze_multiple_images(
                        batch,
                        context=f"Images from DOCX document '{doc_name}'"
                    )
                    if description:
                        image_descriptions.append(description)
        
        # Combine text and image descriptions
        combined_content_parts: List[str] = []
        
        if all_text_parts:
            combined_content_parts.append("\n\n".join(all_text_parts))
        
        if image_descriptions:
            combined_content_parts.append("\n[Image Content]\n" + "\n\n".join(image_descriptions))
        
        if not combined_content_parts:
            logger.info("[DOCX] No extractable content found in '%s'", doc_name)
            return chunks
        
        combined_text = "\n\n".join(combined_content_parts)
        
        doc_metadata = {
            "name": doc_name,
            "page": 1,
            "storage_path": storage_path,
            "has_images": len(all_images) > 0,
        }
        chunks.extend(chunk_document(doc_metadata, combined_text))
        
        return chunks


class XlsxDocumentProcessor(DocumentProcessor):
    """Process XLSX spreadsheets into text chunks.
    
    Treats each worksheet as a separate "page" for chunking purposes.
    Extracts data from cells, preserving row/column structure.
    Extracts and analyzes embedded images using Bedrock vision model.
    """

    def _format_cell_value(self, cell) -> str:
        """Convert cell value to string, handling different types."""
        if cell.value is None:
            return ""
        
        # Handle different cell types
        if cell.is_date and cell.value:
            try:
                return str(cell.value)
            except Exception:
                return str(cell.value)
        
        return str(cell.value).strip()

    def _extract_sheet_text(self, sheet) -> str:
        """Extract all text from a worksheet, preserving structure."""
        rows_text: List[str] = []
        
        for row in sheet.iter_rows():
            cell_values = [self._format_cell_value(cell) for cell in row]
            # Skip completely empty rows
            if not any(cell_values):
                continue
            # Join cells with separator
            rows_text.append(" | ".join(cell_values))
        
        return "\n".join(rows_text)

    def _extract_images_from_workbook(self, file_bytes: bytes) -> Dict[str, List[Tuple[bytes, str]]]:
        """Extract images from workbook, organized by sheet name.
        
        Note: We need to re-open the workbook without read_only mode to access images.
        
        Returns:
            Dict mapping sheet names to lists of (image_bytes, format) tuples
        """
        images_by_sheet: Dict[str, List[Tuple[bytes, str]]] = {}
        
        try:
            # Open workbook without read_only to access images
            workbook = load_workbook(io.BytesIO(file_bytes), data_only=False, read_only=False)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_images: List[Tuple[bytes, str]] = []
                
                # openpyxl stores images in _images attribute
                if hasattr(sheet, '_images'):
                    for image in sheet._images:
                        try:
                            # Get image data
                            image_data = image._data()
                            if image_data:
                                # Try to determine format from the image
                                # Default to png
                                img_format = "png"
                                if hasattr(image, 'path') and image.path:
                                    ext = Path(image.path).suffix.lower()
                                    if ext in ['.jpg', '.jpeg']:
                                        img_format = "jpeg"
                                    elif ext == '.gif':
                                        img_format = "gif"
                                    elif ext == '.webp':
                                        img_format = "webp"
                                
                                sheet_images.append((image_data, img_format))
                        except Exception as e:
                            logger.debug("Failed to extract image from sheet '%s': %s", sheet_name, e)
                
                if sheet_images:
                    images_by_sheet[sheet_name] = sheet_images
            
            workbook.close()
            
        except Exception as e:
            logger.debug("Failed to extract images from workbook: %s", e)
        
        return images_by_sheet

    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        chunks: list[Chunk] = []
        
        # First, extract images from the workbook (separate pass)
        images_by_sheet = self._extract_images_from_workbook(file_bytes)
        
        try:
            # data_only=True gets computed values instead of formulas
            workbook = load_workbook(
                io.BytesIO(file_bytes), 
                data_only=True, 
                read_only=True
            )
        except Exception as e:
            logger.error("[XLSX] Failed to open workbook '%s': %s", doc_name, e)
            raise ValueError(f"Failed to open XLSX workbook '{doc_name}': {e}")
        
        try:
            sheet_names = workbook.sheetnames
            
            if not sheet_names:
                logger.info("[XLSX] No worksheets found in '%s'", doc_name)
                return chunks
            
            for page_number, sheet_name in enumerate(sheet_names, start=1):
                try:
                    sheet = workbook[sheet_name]
                except Exception as e:
                    logger.warning("[XLSX] Failed to access sheet '%s' in '%s': %s", sheet_name, doc_name, e)
                    continue
                
                sheet_text = self._extract_sheet_text(sheet)
                sheet_images = images_by_sheet.get(sheet_name, [])
                
                # Analyze images if present
                image_descriptions: List[str] = []
                if sheet_images:
                    logger.info(
                        "[XLSX] Found %d image(s) in sheet '%s' of '%s'",
                        len(sheet_images), sheet_name, doc_name
                    )
                    
                    if len(sheet_images) == 1:
                        description = self.analyze_image(
                            sheet_images[0][0],
                            sheet_images[0][1],
                            context=f"Image from sheet '{sheet_name}' of Excel workbook '{doc_name}'"
                        )
                        if description:
                            image_descriptions.append(description)
                    else:
                        description = self.analyze_multiple_images(
                            sheet_images,
                            context=f"Images from sheet '{sheet_name}' of Excel workbook '{doc_name}'"
                        )
                        if description:
                            image_descriptions.append(description)
                
                # Combine text and image descriptions
                combined_content_parts: List[str] = []
                
                if sheet_text.strip():
                    # Include sheet name as context for the chunk
                    sheet_header = f"[Sheet: {sheet_name}]\n\n"
                    combined_content_parts.append(sheet_header + sheet_text)
                
                if image_descriptions:
                    combined_content_parts.append("\n[Image Content]\n" + "\n".join(image_descriptions))
                
                if not combined_content_parts:
                    logger.debug("[XLSX] Skipping empty sheet '%s' in '%s'", sheet_name, doc_name)
                    continue
                
                full_text = "\n\n".join(combined_content_parts)
                
                doc_metadata = {
                    "name": doc_name,
                    "page": page_number,
                    "sheet_name": sheet_name,
                    "storage_path": storage_path,
                    "has_images": len(sheet_images) > 0,
                }
                chunks.extend(chunk_document(doc_metadata, full_text))
        
        finally:
            workbook.close()
        
        if not chunks:
            logger.info("[XLSX] No extractable content found in '%s'", doc_name)
        
        return chunks


class PptxDocumentProcessor(DocumentProcessor):
    """Process PPTX presentations into text chunks.
    
    Treats each slide as a separate "page" for chunking purposes.
    Extracts text from shapes, text frames, and tables.
    Extracts and analyzes images using Bedrock vision model.
    """

    def _shape_has_image(self, shape) -> bool:
        """Check if a shape is or contains an image."""
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
            return True
        # Accessing placeholder_format on non-placeholder shapes raises ValueError
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                if shape.placeholder_format.type and 'PICTURE' in str(shape.placeholder_format.type):
                    return True
        except ValueError:
            # Shape is not a placeholder, which is fine
            pass
        return False

    def _extract_image_from_shape(self, shape) -> Optional[Tuple[bytes, str]]:
        """Extract image bytes and format from a picture shape.
        
        Args:
            shape: A python-pptx shape object
            
        Returns:
            Tuple of (image_bytes, format) or None if extraction fails
        """
        try:
            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                # Get the image from the shape
                image = shape.image
                image_bytes = image.blob
                
                # Determine format from content type
                content_type = image.content_type
                img_format = SUPPORTED_IMAGE_FORMATS.get(content_type, "png")
                
                return (image_bytes, img_format)
        except Exception as e:
            logger.debug("Failed to extract image from PPTX shape: %s", e)
        
        return None

    def _get_images_from_slide(self, slide) -> List[Tuple[bytes, str]]:
        """Extract all images from a slide.
        
        Args:
            slide: A python-pptx slide object
            
        Returns:
            List of tuples (image_bytes, format)
        """
        images = []
        
        for shape in slide.shapes:
            if self._shape_has_image(shape):
                img_data = self._extract_image_from_shape(shape)
                if img_data:
                    images.append(img_data)
            
            # Check grouped shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if self._shape_has_image(sub_shape):
                        img_data = self._extract_image_from_shape(sub_shape)
                        if img_data:
                            images.append(img_data)
        
        return images

    def _extract_shape_text(self, shape) -> Optional[str]:
        """Extract text from a shape if it has a text frame."""
        if not shape.has_text_frame:
            return None
        
        text_parts: List[str] = []
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = "".join(run.text for run in paragraph.runs)
            if paragraph_text.strip():
                text_parts.append(paragraph_text.strip())
        
        return "\n".join(text_parts) if text_parts else None

    def _extract_table_text(self, table) -> str:
        """Extract text from a PowerPoint table."""
        rows_text: List[str] = []
        for row in table.rows:
            cell_texts: List[str] = []
            for cell in row.cells:
                cell_text_parts: List[str] = []
                for paragraph in cell.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        cell_text_parts.append(para_text.strip())
                cell_texts.append(" ".join(cell_text_parts))
            rows_text.append(" | ".join(cell_texts))
        return "\n".join(rows_text)

    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        chunks: list[Chunk] = []
        
        try:
            presentation = Presentation(io.BytesIO(file_bytes))
        except Exception as e:
            logger.error("[PPTX] Failed to open presentation '%s': %s", doc_name, e)
            raise ValueError(f"Failed to open PPTX presentation '{doc_name}': {e}")
        
        if not presentation.slides:
            logger.info("[PPTX] No slides found in '%s'", doc_name)
            return chunks
        
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text_parts: List[str] = []
            slide_images: List[Tuple[bytes, str]] = []
            
            # Extract slide title if available
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = self._extract_shape_text(slide.shapes.title)
                if title_text:
                    slide_text_parts.append(f"# {title_text}")
            
            # Extract text and images from all shapes
            for shape in slide.shapes:
                # Skip title as we already processed it
                if shape == slide.shapes.title:
                    continue
                
                # Extract images
                if self._shape_has_image(shape):
                    img_data = self._extract_image_from_shape(shape)
                    if img_data:
                        slide_images.append(img_data)
                    continue
                
                # Handle tables
                if shape.has_table:
                    table_text = self._extract_table_text(shape.table)
                    if table_text.strip():
                        slide_text_parts.append(table_text)
                    continue
                
                # Handle grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub_shape in shape.shapes:
                        if self._shape_has_image(sub_shape):
                            img_data = self._extract_image_from_shape(sub_shape)
                            if img_data:
                                slide_images.append(img_data)
                        elif sub_shape.has_text_frame:
                            text = self._extract_shape_text(sub_shape)
                            if text:
                                slide_text_parts.append(text)
                    continue
                
                # Handle regular shapes with text
                if shape.has_text_frame:
                    text = self._extract_shape_text(shape)
                    if text:
                        slide_text_parts.append(text)
            
            # Analyze images if present
            image_descriptions: List[str] = []
            if slide_images:
                logger.info(
                    "[PPTX] Found %d image(s) on slide %d of '%s'",
                    len(slide_images), slide_number, doc_name
                )
                
                if len(slide_images) == 1:
                    description = self.analyze_image(
                        slide_images[0][0],
                        slide_images[0][1],
                        context=f"Image from slide {slide_number} of PowerPoint '{doc_name}'"
                    )
                    if description:
                        image_descriptions.append(description)
                else:
                    description = self.analyze_multiple_images(
                        slide_images,
                        context=f"Images from slide {slide_number} of PowerPoint '{doc_name}'"
                    )
                    if description:
                        image_descriptions.append(description)
            
            # Combine text and image descriptions
            combined_content_parts: List[str] = []
            
            if slide_text_parts:
                combined_content_parts.append("\n\n".join(slide_text_parts))
            
            if image_descriptions:
                combined_content_parts.append("\n[Image Content]\n" + "\n".join(image_descriptions))
            
            if not combined_content_parts:
                logger.debug("[PPTX] Slide %d of '%s' has no extractable content", slide_number, doc_name)
                continue
            
            slide_text = "\n\n".join(combined_content_parts)
            
            doc_metadata = {
                "name": doc_name,
                "page": slide_number,
                "storage_path": storage_path,
                "has_images": len(slide_images) > 0,
            }
            chunks.extend(chunk_document(doc_metadata, slide_text))
        
        if not chunks:
            logger.info("[PPTX] No extractable content found in '%s'", doc_name)
        
        return chunks


class XmlDocumentProcessor(DocumentProcessor):
    """Process XML documents into text chunks.
    
    Extracts text content from XML elements, handling different XML structures.
    Preserves element hierarchy as context where meaningful.
    """

    def _clean_tag(self, tag: str) -> str:
        """Remove namespace prefix from XML tag."""
        # Handle {namespace}tag format
        if '}' in tag:
            return tag.split('}', 1)[1]
        return tag

    def _extract_text_recursive(
        self, 
        element: ET.Element, 
        depth: int = 0,
        max_depth: int = 50
    ) -> List[str]:
        """Recursively extract text from XML elements."""
        if depth > max_depth:
            return []
        
        text_parts: List[str] = []
        
        # Get element's direct text
        if element.text and element.text.strip():
            text_parts.append(element.text.strip())
        
        # Process child elements
        for child in element:
            child_texts = self._extract_text_recursive(child, depth + 1, max_depth)
            text_parts.extend(child_texts)
            
            # Get tail text (text after child element)
            if child.tail and child.tail.strip():
                text_parts.append(child.tail.strip())
        
        return text_parts

    def _extract_structured_text(self, root: ET.Element) -> str:
        """Extract text while preserving some structure from XML."""
        # Common document-like XML structures
        document_tags = {'document', 'doc', 'article', 'content', 'body', 'text'}
        section_tags = {'section', 'chapter', 'part', 'div', 'paragraph', 'p', 'para'}
        
        output_parts: List[str] = []
        
        def process_element(element: ET.Element, level: int = 0) -> None:
            tag = self._clean_tag(element.tag).lower()
            
            # Check if this is a section-like element
            is_section = tag in section_tags
            
            element_text_parts: List[str] = []
            
            # Get direct text
            if element.text and element.text.strip():
                element_text_parts.append(element.text.strip())
            
            # Process children
            for child in element:
                process_element(child, level + 1)
                # Get tail text
                if child.tail and child.tail.strip():
                    element_text_parts.append(child.tail.strip())
            
            # Add text from this element
            if element_text_parts:
                element_text = " ".join(element_text_parts)
                if is_section:
                    output_parts.append("\n" + element_text)
                elif element_text not in " ".join(output_parts):
                    # Avoid duplicating text already added by children
                    pass
        
        # Try structured extraction first
        process_element(root)
        
        if output_parts:
            return "\n\n".join(part.strip() for part in output_parts if part.strip())
        
        # Fallback to simple recursive extraction
        all_texts = self._extract_text_recursive(root)
        return "\n".join(all_texts)

    def process(
        self,
        file_bytes: bytes,
        *,
        chunk_document: Chunker,
        storage_path: str,
        doc_name: str,
    ) -> list[Chunk]:
        chunks: list[Chunk] = []
        
        # Decode XML content
        try:
            xml_content = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            try:
                xml_content = file_bytes.decode('latin-1')
            except UnicodeDecodeError:
                logger.error("[XML] Failed to decode '%s' with UTF-8 or Latin-1", doc_name)
                raise ValueError(f"Failed to decode XML document '{doc_name}'")
        
        # Parse XML
        try:
            root = ET.fromstring(xml_content)
        except ET.ParseError as e:
            logger.error("[XML] Failed to parse XML document '%s': %s", doc_name, e)
            raise ValueError(f"Failed to parse XML document '{doc_name}': {e}")
        
        # Extract text content
        extracted_text = self._extract_structured_text(root)
        
        if not extracted_text.strip():
            # Fallback: try simple text extraction
            all_texts = self._extract_text_recursive(root)
            extracted_text = "\n".join(all_texts)
        
        if not extracted_text.strip():
            logger.info("[XML] No text content found in '%s'", doc_name)
            return chunks
        
        # Add root element info as context
        root_tag = self._clean_tag(root.tag)
        
        doc_metadata = {
            "name": doc_name,
            "page": 1,
            "root_element": root_tag,
            "storage_path": storage_path,
        }
        chunks.extend(chunk_document(doc_metadata, extracted_text))
        
        return chunks


_CONTENT_TYPE_PROCESSORS: Dict[str, Type[DocumentProcessor]] = {
    "application/pdf": PDFDocumentProcessor,
    "text/plain": TextDocumentProcessor,
    "text/csv": TextDocumentProcessor,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocxDocumentProcessor,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": XlsxDocumentProcessor,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": PptxDocumentProcessor,
    "application/xml": XmlDocumentProcessor,
    "text/xml": XmlDocumentProcessor,
}

_EXTENSION_FALLBACKS: Dict[str, Type[DocumentProcessor]] = {
    ".pdf": PDFDocumentProcessor,
    ".txt": TextDocumentProcessor,
    ".csv": TextDocumentProcessor,
    ".docx": DocxDocumentProcessor,
    ".xlsx": XlsxDocumentProcessor,
    ".pptx": PptxDocumentProcessor,
    ".xml": XmlDocumentProcessor,
}


def get_document_processor(content_type: Optional[str], storage_path: str) -> DocumentProcessor:
    """Return a processor that matches the given content type or file extension."""

    processor_cls: Optional[Type[DocumentProcessor]] = None

    if content_type:
        processor_cls = _CONTENT_TYPE_PROCESSORS.get(content_type.lower())

    if processor_cls is None:
        extension = Path(storage_path).suffix.lower()
        processor_cls = _EXTENSION_FALLBACKS.get(extension)

    if processor_cls is None:
        raise ValueError(f"No document processor available for type '{content_type}' or extension '{storage_path}'.")

    return processor_cls()
